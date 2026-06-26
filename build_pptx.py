"""Build the local-LLM deck natively into the Data Society pptx template.

Opens the brand template (keeps its masters/layouts/theme/logo), clears the 4 sample
slides, then emits our deck using the template's own layouts so everything inherits the
Data Society fonts/colors. Tables are styled with the brand palette; the two side-by-side
animations embed as gifs (animate in Google Slides; static first-frame in PowerPoint).

Run: python3 build_pptx.py
Out: local-llm-datasociety.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

TPL = "/Users/john.connolly/Downloads/Data Society slide template.pptx"
REPO = Path("/Users/john.connolly/tmp-projects/local-llm")
OUT = REPO / "local-llm-datasociety.pptx"

# brand palette (from theme1.xml)
DK     = RGBColor(0x25, 0x11, 0x44)
PURPLE = RGBColor(0x83, 0x5B, 0xA3)
DK2    = RGBColor(0x41, 0x2B, 0x71)
RED    = RGBColor(0xFF, 0x35, 0x0E)
YELLOW = RGBColor(0xFF, 0xED, 0x4C)
LIGHT  = RGBColor(0xF1, 0xF1, 0xF4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation(TPL)
SW, SH = prs.slide_width, prs.slide_height
LAY = {l.name: l for l in prs.slide_masters[0].slide_layouts}        # white content master
PURPLE_LAY = {l.name: l for l in prs.slide_masters[1].slide_layouts}  # purple section master (#251144)

# clear the 4 sample slides, drop the relationship too, else the orphaned slide
# parts stay in the package and collide with new slideN.xml names (duplicate zip entries).
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def _md_runs(p, text):
    """Add runs to paragraph p, honoring **bold** and `code` markers."""
    import re
    parts = re.split(r"(\*\*.+?\*\*|`.+?`)", text)
    for part in parts:
        if not part:
            continue
        r = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            r.text = part[2:-2]; r.font.bold = True
        elif part.startswith("`") and part.endswith("`"):
            r.text = part[1:-1]; r.font.name = "Consolas"
        else:
            r.text = part


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _no_autofit(tf):
    """Kill TEXT_TO_FIT_SHAPE autofit so every app renders the same explicit size."""
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def _style_title(ph, size=30):
    _no_autofit(ph.text_frame)
    for p in ph.text_frame.paragraphs:
        p.font.size = Pt(size)
        for r in p.runs:
            r.font.size = Pt(size)


def _fmt_body(tf, lvl_sizes={0: 18, 1: 15}):
    _no_autofit(tf)
    for p in tf.paragraphs:
        sz = Pt(lvl_sizes.get(p.level, 15))
        p.font.size = sz
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = sz


def title_slide(title, subtitle, byline="", tldr="", tldr_gif="", notes=""):
    s = prs.slides.add_slide(LAY["Title Slide No Image"])
    def place(ph, left, top, w, h):
        ph.left, ph.top, ph.width, ph.height = Inches(left), Inches(top), Inches(w), Inches(h)
    # title
    s.placeholders[0].text = title
    place(s.placeholders[0], 0.9, 1.9, 10.2, 1.5)
    _style_title(s.placeholders[0], size=36)
    # snappy one-line subtitle
    if len(s.placeholders) > 1:
        s.placeholders[1].text = subtitle
        place(s.placeholders[1], 0.9, 3.4, 10.2, 0.6)
        tf = s.placeholders[1].text_frame
        _no_autofit(tf)
        for p in tf.paragraphs:
            p.font.size = Pt(20); p.space_after = Pt(0)
            for r in p.runs:
                r.font.size = Pt(20); r.font.color.rgb = DK2
    # TL;DR as a separate styled callout box (brand tint, red label, bold verdicts)
    if tldr:
        box_top, box_h = (4.05, 1.55) if tldr_gif else (4.2, 1.15)
        box_w = 8.2 if tldr_gif else 10.2
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.9), Inches(box_top), Inches(box_w), Inches(box_h))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF1, 0xEC, 0xF7)
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame; _no_autofit(tf)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3); tf.margin_right = Inches(2.6 if tldr_gif else 0.3)
        p = tf.paragraphs[0]
        lab = p.add_run(); lab.text = "TL;DR   "
        lab.font.bold = True; lab.font.size = Pt(19); lab.font.color.rgb = RED
        _md_runs(p, tldr)
        for r in p.runs[1:]:
            r.font.size = Pt(18); r.font.color.rgb = DK
        if tldr_gif:  # inside the box, right-aligned, vertically centered
            gh = 1.25
            gp = s.shapes.add_picture(str(tldr_gif), 0, 0, height=Inches(gh))
            gp.left = int(Inches(0.9 + box_w - 0.22) - gp.width)
            gp.top = int(Inches(box_top + (box_h - gh) / 2))
    # byline: bottom-RIGHT (opposite the bottom-left footer logo), right-aligned, small.
    # multi-line via \n.
    if byline:
        tb = s.shapes.add_textbox(Inches(4.6), Inches(5.7), Inches(6.5), Inches(0.85))
        tf = tb.text_frame; _no_autofit(tf)
        for i, line in enumerate(byline.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = line
            r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DK2
    _notes(s, notes)
    return s


def _tldr_two_gif(s):
    """Custom title-slide TL;DR: two reaction gifs, each with a parenthetical caption."""
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.9), Inches(3.9), Inches(6.5), Inches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF1, 0xEC, 0xF7)
    box.line.fill.background(); box.shadow.inherit = False

    def tb(text, left, top, w, sz, color, italic=False, bold=False):
        b = s.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(0.35))
        _no_autofit(b.text_frame)
        p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.color.rgb = color
        r.font.italic = italic; r.font.bold = bold

    tb("TL;DR", 1.0, 4.55, 1.0, 19, RED, bold=True)
    # column 1: bounded coding -> "not bad"
    tb("bounded coding", 2.1, 3.98, 2.4, 14, DK)
    g1 = s.shapes.add_picture(str(REPO / "viz" / "not-bad-slow.gif"), 0, 0, height=Inches(0.82))
    g1.left = int(Inches(3.3) - g1.width // 2); g1.top = int(Inches(4.3))
    tb("functionally equivalent, 10x slower", 1.9, 5.18, 2.8, 11, DK2, italic=True)
    # column 2: open-ended repos -> conceited
    tb("Open-ended repos:", 4.6, 3.98, 2.4, 14, DK)
    g2 = s.shapes.add_picture(str(REPO / "viz" / "conceited.gif"), 0, 0, height=Inches(0.82))
    g2.left = int(Inches(5.8) - g2.width // 2); g2.top = int(Inches(4.3))
    tb("(not quite)", 4.6, 5.18, 2.4, 12, DK2, italic=True)


def bullets_slide(title, items, notes="", font=18):
    """items: list of (text, level)."""
    s = prs.slides.add_slide(LAY["OBJECT"])  # white content layout WITH the DS logo
    s.placeholders[0].text = title
    _style_title(s.placeholders[0])
    tf = s.placeholders[1].text_frame
    tf.clear()
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        _md_runs(p, text)
    _fmt_body(tf, {0: font, 1: font - 3})
    _notes(s, notes)
    return s


def _title_only(title):
    s = prs.slides.add_slide(LAY["TITLE_ONLY"])  # white title-only layout WITH the DS logo
    s.placeholders[0].text = title
    _style_title(s.placeholders[0])
    return s


def table_slide(title, headers, rows, caption="", notes="", col_w=None,
                top=1.7, font=13):
    s = _title_only(title)
    nrow, ncol = len(rows) + 1, len(headers)
    width = SW - Inches(1.2)
    left = Inches(0.6)
    height = Inches(0.5) * nrow
    gfx = s.shapes.add_table(nrow, ncol, left, Inches(top), width, height)
    tbl = gfx.table
    if col_w:
        total = sum(col_w)
        for ci, w in enumerate(col_w):
            tbl.columns[ci].width = int((SW - Inches(1.2)) * w / total)
    # header
    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = PURPLE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame; tf.clear(); p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(font + 1)
        r.font.name = "Arial"
    # body
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.clear(); p = tf.paragraphs[0]
            bold = val.startswith("**") and val.endswith("**")
            txt = val[2:-2] if bold else val
            r = p.add_run(); r.text = txt
            r.font.size = Pt(font); r.font.color.rgb = DK; r.font.name = "Arial"
            r.font.bold = bold
    if caption:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(top) + height + Inches(0.2),
                                  SW - Inches(1.2), Inches(1.1))
        _no_autofit(tb.text_frame)
        p = tb.text_frame.paragraphs[0]
        _md_runs(p, caption)
        for r in p.runs:
            r.font.size = Pt(13); r.font.color.rgb = DK2
    _notes(s, notes)
    return s


def image_slide(title, img, notes="", width_in=9.0, sub=""):
    s = _title_only(title)
    pic = s.shapes.add_picture(str(img), 0, 0, width=Inches(width_in))
    # clamp so it never overflows the content band (1.55in down to ~6.5in)
    max_h = Inches(4.55)
    if pic.height > max_h:
        scale = max_h / pic.height
        pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
    pic.left = int((SW - pic.width) / 2)
    pic.top = Inches(1.55)
    if sub:
        sub_top = pic.top + pic.height + Inches(0.06)
        tb = s.shapes.add_textbox(Inches(0.6), sub_top, SW - Inches(1.2), Inches(0.4))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _md_runs(p, sub)
        for r in p.runs:
            r.font.size = Pt(13); r.font.color.rgb = DK2
    _notes(s, notes)
    return s


def section_slide(kicker, title, notes=""):
    """Full-purple section divider (uses the purple master) with a yellow kicker + white title.
    Weaves the template's white/purple/white rhythm between acts."""
    lay = PURPLE_LAY.get("TITLE_ONLY") or PURPLE_LAY.get("Title Only - No Logo") \
        or list(PURPLE_LAY.values())[1]
    s = prs.slides.add_slide(lay)
    kb = s.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(10.2), Inches(0.6))
    _no_autofit(kb.text_frame)
    kr = kb.text_frame.paragraphs[0].add_run()
    kr.text = kicker; kr.font.size = Pt(20); kr.font.bold = True; kr.font.color.rgb = YELLOW
    tb = s.shapes.add_textbox(Inches(0.9), Inches(3.15), Inches(10.2), Inches(1.9))
    _no_autofit(tb.text_frame)
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = title; tr.font.size = Pt(40); tr.font.bold = True; tr.font.color.rgb = WHITE
    _notes(s, notes)
    return s


B = lambda t: (t, 0)   # level-0 bullet
B1 = lambda t: (t, 1)  # level-1 bullet

# ============================ THE DECK ============================
_title = title_slide(
    "Can a local LLM replace cloud Claude Code?",
    "Three weeks, three machines: from a Raspberry Pi to a gray-market used Mac Studio.",
    byline="John Connolly, Lead Product Engineer & tinkerer\nJune 2026",
    notes="[0:30]\nSAY: \"Three weeks ago I asked a simple question: how much LLM work can I do on my own hardware? I spent the price of a used Mac Studio finding out. Could I get a toy conversational model running? How far could I push it? The journey is the talk: for bounded coding, a local LLM can functionally tie the frontier, and it's free; for open-ended repo work, spoiler, you still want cloud.\"\n\nCUE: Open cold, don't read the title. Point at the TL;DR box on 'bounded' then 'open-ended.' A measurement talk, not a vibes talk.")
# --- slide 2: the "No." gag (deadpan answer to the title question) ---
_joke = prs.slides.add_slide(LAY["TITLE_ONLY"])  # keep the DS logo
_joke.placeholders[0].text = ""                  # no title
_jnb = _joke.shapes.add_textbox(Inches(0), Inches(1.0), SW, Inches(2.5))
_no_autofit(_jnb.text_frame); _jnp = _jnb.text_frame.paragraphs[0]; _jnp.alignment = PP_ALIGN.CENTER
_jnr = _jnp.add_run(); _jnr.text = "No."
_jnr.font.size = Pt(180); _jnr.font.bold = True; _jnr.font.color.rgb = PURPLE
_jrp = _joke.shapes.add_picture(str(REPO / "viz" / "rageface.png"), 0, 0, height=Inches(2.0))
_jrp.left = int((SW - _jrp.width) // 2); _jrp.top = Inches(3.65)
_jsb = _joke.shapes.add_textbox(Inches(0), Inches(6.0), SW, Inches(0.7))
_no_autofit(_jsb.text_frame); _jsp = _jsb.text_frame.paragraphs[0]; _jsp.alignment = PP_ALIGN.CENTER
_jsr = _jsp.add_run(); _jsr.text = "Thank you for coming to my Ted Talk."
_jsr.font.size = Pt(26); _jsr.font.italic = True; _jsr.font.color.rgb = DK2
_notes(_joke, "[0:15]\nSAY: (deadpan, pause) \"...No. Thank you for coming to my Ted Talk.\" (beat, let it land, then break) \"...okay. The honest answer is more interesting than that. It's 'it depends,' and the where-it-depends is the whole talk. Let's go.\"\n\nCUE: Deliver totally straight. Long pause after 'No.' Pretend you're wrapping up, then smile and move on. Resets the room, buys goodwill for the next 30 minutes.")

# --- slide 3: "No, but seriously..." -> the real TL;DR ---
_seriously = _title_only("No, but seriously…")
_tldr_two_gif(_seriously)
_notes(_seriously, "[0:40]\nSAY: \"Okay. Seriously. Here's the honest one-slide answer, and the rest of the talk is me earning it. For bounded coding, write-this-function, fix-this-bug, a local model on my own hardware is functionally equivalent to frontier cloud, it just runs about ten times slower, and it's free. For open-ended repo work, multi-file, vague bug report, sprawling context, it's not quite there, cloud still wins. Two regimes, one line each. Now let me show you how I know.\"\n\nCUE: This is the thesis. Land it slowly, point bounded then open-ended. Everything after this slide is evidence for these two lines.")

_who = bullets_slide("whoami",
    [B("**Lead Product Engineer @ Data Society**, since Jan 2025"),
     B("Day-to-day: rapid prototypes, enabling engineering across the org, plugging the gaps (internal CMS & co.)"),
     B("CS degree; member of Stony Brook's **Group for Logic & Formal Semantics**; contributed to a paper in MIT's *Artificial Life* journal"),
     B("Shipping software since **2007**. Learn by doing, not reading spec sheets (which is how this talk happened)"),
     B("My personal projects include: OpenWRT home network, home automation, a personal-finance app, Raspberry Pi projects with my daughter"),
     B("Not-programming: sourdough from scratch, surfing when the swell allows, and a ~10-acre subsistence farm someday (Adirondacks / Catskills / coastal Maine)"),
     B("Recently moved my own coding from **Cursor to Claude Code** (receipts below)")],
    "[0:40]\nSAY: \"Quick bit about me. Lead Product Engineer at Data Society, started January of last year. Day to day I build prototypes, enable engineering across the org, and plug gaps wherever they show up, a lot of that on internal CMS. CS degree, did time in Stony Brook's logic and formal semantics group, even have a paper in MIT's Artificial Life journal from way back. Shipping software since 2007, I learn by getting my hands on things, which is how this talk happened. Off the clock it's home networking, home automation, sourdough, surfing, and someday a little subsistence farm. And I recently moved my own coding from Cursor to Claude Code, part of why I wanted to know if I could run it at home. That heatmap at the bottom is my last year of dev activity, three sources stacked. It's incomplete but directionally honest.\"\n\nCUE: 40 seconds, don't over-share. Point at the heatmap on 'incomplete but directionally honest.' Sets up the whole talk.",
    font=14)
_wp = _who.shapes.add_picture(str(REPO / "viz" / "activity.png"), 0, 0, width=Inches(7.4))
_wp.left = int((SW - _wp.width) // 2); _wp.top = Inches(4.95)
_wc = _who.shapes.add_textbox(Inches(0.6), _wp.top + _wp.height + Inches(0.05), SW - Inches(1.2), Inches(0.55))
_no_autofit(_wc.text_frame); _wcp = _wc.text_frame.paragraphs[0]; _wcp.alignment = PP_ALIGN.CENTER
_wcr = _wcp.add_run()
_wcr.text = "My dev activity, last 52 weeks (Claude Code / Cursor / Bitbucket, stacked). An incomplete but directionally accurate snapshot - Claude Code keeps only ~6 weeks of logs, so its band is thin on the left."
_wcr.font.size = Pt(11); _wcr.font.color.rgb = DK2

bullets_slide("Agenda, where we're going (~30m)",
    [B("**The route:** Raspberry Pi, Air, Studio"),
     B("**Reality check:** you can't buy parity, the hardware call"),
     B("**The measured verdict:** one-shot, agent loop, the build test"),
     B("**Economics, recommendation, honest caveats**")],
    "[0:30]\nSAY: \"Four parts. Part one is the route across three machines. Part two, the uncomfortable truth that you can't buy your way to parity. Part three is the heart, the benchmarks, and it's the part that changed my mind. Part four, the money and the recommendation. I'll leave time for questions.\"\n\nCUE: Don't dwell. If running long, Part 1 is the part to compress.")

section_slide("PART ONE", "The route: three machines",
    notes="[0:05]\nSAY: \"Part one, the route. How I got from a thirty-five-dollar Raspberry Pi to a gray-market Mac Studio, and what each machine taught me.\"\n\nCUE: Quick beat, don't linger.")

bullets_slide("The question: how capable could a home LLM get?",
    [B("\"I want to run an LLM at home. How capable could it actually be?\""),
     B("Shoot for the moon! Good enough to be my Claude Code backend, so I can stop paying for cloud."),
     B("Two things to find out:"),
     B1("Is it possible, and at what hardware/cost?"),
     B1("Is it good enough? Measured, not vibes.")],
    "[1:00]\nSAY: \"The real question was simple, and a little greedy: I want to run an LLM at home, how capable could it actually get? The bar: good enough to be my Claude Code backend, so I can stop paying for cloud. That splits in two, is it possible, and at what hardware and cost, and is it good enough, measured, not vibes.\"\n\nCUE: Lead with motivation, not a spec sheet. Plant that 'good enough' itself splits later (bounded vs open-ended) with opposite answers. Last bullet is the north star.")

_pi5 = bullets_slide("Why start on a Raspberry Pi? Why not?",
    [B("Learn by doing, not reading spec sheets"),
     B("Already had a Pi (projects with my daughter)"),
     B("Ran a small LLM at home. Could it do real dev work?")],
    "[1:00]\nSAY: \"I learn by jumping in, not reading a spec sheet. The origin's mundane: I had a Raspberry Pi from tinkering projects with my daughter, and wondered if it could run a little LLM for the house. It did, so I got greedy and asked the real question, could it handle my actual dev workload?\"\n\nCUE: The human hook, keeps the Pi from looking naive. When it turns out a dead end, that was intentional, I wanted to feel the limits, not predict them.")
_pi5.placeholders[1].width = Inches(6.3)
_pi5p = _pi5.shapes.add_picture(str(REPO / "viz" / "pi5-aihat.png"), Inches(7.0), Inches(2.45), width=Inches(3.45))
_pi5c = _pi5.shapes.add_textbox(Inches(6.85), _pi5p.top + _pi5p.height + Inches(0.12), Inches(3.75), Inches(0.6))
_no_autofit(_pi5c.text_frame)
_pi5cp = _pi5c.text_frame.paragraphs[0]; _pi5cp.alignment = PP_ALIGN.CENTER
_pi5cr = _pi5cp.add_run(); _pi5cr.text = "Raspberry Pi 5 + AI HAT+ 2 (bought from Vilros)"
_pi5cr.font.size = Pt(12); _pi5cr.font.italic = True; _pi5cr.font.color.rgb = DK2

table_slide("Bounded coding vs open-ended repos",
    ["", "Bounded coding", "Open-ended repos"],
    [["Scope", "one function, script, or algorithm; a file or a few you already know", "a vague bug or feature across a large, unfamiliar codebase"],
     ["Context", "small, fits in your head", "huge (64k+), must be discovered first"],
     ["Agent loop", "few turns, converges fast", "many turns, sustained reasoning"],
     ["Examples", "a parser, a failing unit test, a data structure, a CLI tool, LeetCode", "refactor across the repo, a feature touching 12 files, real SWE-bench"],
     ["Benchmark", "**HumanEval, LiveCodeBench** (function-level)", "**SWE-bench Verified** (repository-level)"]],
    "**Which side does local actually win? Where does cloud stay ahead?**",
    "[1:00]\nSAY: \"Two genuinely different jobs. Bounded coding is a self-contained problem you hold in your head and finish in a few turns, a parser, a failing test, a small CLI tool, measured by HumanEval and LiveCodeBench. Open-ended is a vague task across a big unfamiliar codebase, sixty-thousand-plus tokens, many turns, measured by SWE-bench Verified. So which side does local actually win, and where does cloud stay ahead?\"\n\nCUE: The whole talk hinges on this. Do NOT answer who-wins yet, that's the payoff. Plant the two-axis framing and let the question hang.",
    col_w=[1.2, 4.0, 4.4], font=12)

_q = _title_only("Where does the line fall? (the map we'll test)")
_qp = _q.shapes.add_picture(str(REPO / "viz" / "quadrant.png"), 0, 0, height=Inches(4.35))
_qp.left = int((SW - _qp.width) // 2); _qp.top = Inches(1.4)
_qcap = _q.shapes.add_textbox(Inches(1.7), Inches(5.95), SW - Inches(2.4), Inches(0.7))
_no_autofit(_qcap.text_frame); _qcp = _qcap.text_frame.paragraphs[0]; _qcp.alignment = PP_ALIGN.CENTER
_qcr = _qcp.add_run()
_qcr.text = "Teal circle = my coding tasks, manually classified: they lean open-ended (~40% local / ~60% cloud). Most of my Claude Code use isn't coding at all (ops, Jira/Slack, infra, personal projects) - full 2,177-prompt breakdown next."
_qcr.font.size = Pt(12); _qcr.font.color.rgb = DK2
_notes(_q, "[1:30]\nSAY: \"I built this chart, then looked at my actual history, 2,177 prompts across 68 projects. The teal circle is my coding tasks: they straddle the line and lean cloud, about sixty percent open-ended, exactly where local is weakest. The bigger surprise, most of what I use Claude Code for isn't even coding, it's ops, Jira and Slack, infra, a finance app, building this deck.\"\n\nCUE: The teal split is a manual estimate; the full breakdown is next slide. Land it: my Claude Code is more general agent than code generator.")

_u = _title_only("Half work, half personal (outside work)")
_up = _u.shapes.add_picture(str(REPO / "viz" / "usage.png"), 0, 0, height=Inches(4.95))
_up.left = int((SW - _up.width) // 2); _up.top = Inches(1.45)
_notes(_u, "[1:30]\nSAY: \"The full audit, all 2,177 prompts, classified by project. Two headlines. One, it's almost exactly half DSG work, half personal projects outside work. Two, a huge chunk of the personal half is privacy-sensitive, my finances, my home network, my thermostat, and that's local's natural home, the data never leaves the house. And it's not a compromise, when a local model categorized my credit-card charges it was genuinely good, I didn't need the cloud.\"\n\nCUE: This is why local matters beyond coding, the private bounded-classification work is half my usage.")

_de1 = bullets_slide("Dead end #1: the Pi (tilting at windmills)",
    [B("Pi 5 + Hailo-10H, 40-TOPS NPU. A real gen-AI accelerator."),
     B("It does run LLMs, Hailo GenAI zoo + an Ollama/OpenAI-compatible runtime.¹²"),
     B("But the zoo is **1-3B**, too small for a 64k-context coding agent.³"),
     B("Measured **~7 tok/s**; reviewers find the **CPU often beats it**⁴⁵ (Hailo markets 30-50)."),
     B("**Takeaway:** a capability problem, not a connection one. Wrong job for a coding agent.")],
    "[1:30]\nSAY: \"I verified this on the actual board, and to be fair, the Hailo does run LLMs now, a generative-AI model zoo plus an Ollama-compatible runtime, so you can point Claude Code at it. So why a dead end? Two reasons, both capability, not connection. One, every model in that zoo is one to three billion parameters, far too small for a sixty-four-thousand-token coding agent. Two, it's slow, I measured about seven tokens a second, and reviewers find the plain CPU often beats it, one literally called it more like a decelerator than an accelerator. Great low-power chatbot, never a coding backend.\"\n\nCUE: Corrected since I first built this, own that. Let the gif play. Sources on the slide.",
    font=16)
_de1.placeholders[1].width = Inches(6.7)
_dgif = _de1.shapes.add_picture(str(REPO / "viz" / "fast-furious.gif"), Inches(7.4), Inches(2.2), width=Inches(3.7))
_dcap = _de1.shapes.add_textbox(Inches(7.4), Inches(2.2) + _dgif.height + Inches(0.06), Inches(3.8), Inches(0.8))
_no_autofit(_dcap.text_frame)
_dcp = _dcap.text_frame.paragraphs[0]; _dcp.alignment = PP_ALIGN.CENTER
_dcr = _dcp.add_run(); _dcr.text = "“more like an AI decelerator than an AI accelerator”"
_dcr.font.size = Pt(12); _dcr.font.italic = True; _dcr.font.color.rgb = DK2
_dcp2 = _dcap.text_frame.add_paragraph(); _dcp2.alignment = PP_ALIGN.CENTER
_dcr2 = _dcp2.add_run(); _dcr2.text = "— CNX Software, AI HAT+ 2 review"
_dcr2.font.size = Pt(10); _dcr2.font.color.rgb = DK2
_de1f = _de1.shapes.add_textbox(Inches(0.6), Inches(6.18), SW - Inches(1.2), Inches(0.8))
_no_autofit(_de1f.text_frame); _de1f.text_frame.word_wrap = True
_de1p = _de1f.text_frame.paragraphs[0]
_SRCS = [
    ("1", "Hailo: On-device GenAI on the Pi AI HAT+ 2", "https://hailo.ai/blog/bringing-on-device-generative-ai-to-the-pi-when-and-why-youll-need-the-raspberry-pi-ai-hat-2/"),
    ("2", "Hailo GenAI Model Zoo (GitHub)", "https://github.com/hailo-ai/hailo_model_zoo_genai"),
    ("3", "raspberry.tips: AI HAT+ 2 local LLMs", "https://raspberry.tips/en/raspberrypi-tutorials/raspberry-pi-ai-hat-2-hailo-10h-40-tops-local-llms"),
    ("4", "CNX Software: AI HAT+ 2 review", "https://www.cnx-software.com/2026/01/20/raspberry-pi-ai-hat-2-review-a-40-tops-ai-accelerator-tested-with-computer-vision-llm-and-vlm-workloads/"),
    ("5", "hardware-corner.net: Local LLMs on the AI HAT+ 2", "https://www.hardware-corner.net/local-llms-raspberry-pi-ai-hat-plus-2/"),
]
for _i, (_n, _label, _url) in enumerate(_SRCS):
    _nr = _de1p.add_run(); _nr.text = f"{_n} "
    _nr.font.size = Pt(9); _nr.font.bold = True; _nr.font.color.rgb = DK2
    _lr = _de1p.add_run(); _lr.text = _label
    _lr.font.size = Pt(9); _lr.font.italic = True; _lr.font.color.rgb = PURPLE
    _lr.hyperlink.address = _url
    if _i < len(_SRCS) - 1:
        _sp = _de1p.add_run(); _sp.text = "    "; _sp.font.size = Pt(9)

bullets_slide("Why memory bandwidth is the ceiling (one bit of theory)",
    [B("Generating a token = **read every active weight from memory, once**. The arithmetic is trivial; the *reading* is the job."),
     B("Speed limit, full stop: **tokens/sec ≈ memory bandwidth ÷ bytes read per token**."),
     B("At batch 1 the chip is **starved, not busy** - idling on the memory bus (memory-bound side of the roofline). More teraflops buy nothing."),
     B("So the Pi's **40 TOPS was never the problem** - TOPS is compute; the bus is the wall."),
     B("Reading a 30B model (q8, ~30GB): Pi ~17 GB/s -> **<1 tok/s**; Mac Studio **819 GB/s -> ~27 tok/s**. Same model, ~48x the bus."),
     B("**MoE escape hatch:** activate only ~3B of 30B per token -> bytes-read drops ~10x -> fast *and* smart. (Real numbers in Part 3.)")],
    "[1:00]\nSAY: \"One bit of theory, because it explains every machine in this talk. When a model generates a token, it reads every active weight out of memory, once. The math is trivial; the reading is the whole job. So the speed limit is almost embarrassingly simple: tokens per second is roughly memory bandwidth divided by how many bytes you read per token. At batch one the chip isn't busy, it's starved, sitting idle waiting on the memory bus. That's the punchline for Part One: the Pi's forty TOPS was never the problem, because TOPS is compute, and the wall is the bus. Reading a thirty-gig model, the Pi tops out under one token a second; the Studio I end up buying does about twenty-seven, same model, forty-eight times the bandwidth. The escape hatch is mixture-of-experts: activate only a few billion of the thirty billion per token, read less, go faster without getting dumber. Real numbers in Part Three.\"\n\nCUE: Load-bearing concept slide for the hardware arc. Say it once, slowly. Point at the two numbers, don't read them. Land 'bandwidth, not teraflops.' Foreshadow MoE, don't explain it yet.",
    font=16)

bullets_slide("Anatomy of one token (30B, q4, on the Studio)",
    [B("**Setup:** 30B params at 4-bit ≈ **15GB** of weights in unified memory; ~48 transformer layers. One token = walk all 48, once."),
     B("**Each layer - Attention:** multiply the current token-vector by the Q/K/V/O weight matrices."),
     B1("Memory: stream those weights once; read the KV cache (every past token) + append this token's K/V. Cores: a few multiply-adds, then idle."),
     B("**Each layer - FFN:** multiply by the two big matrices (~2/3 of the layer's weights)."),
     B1("Memory: stream once. Cores: multiply-add, then idle."),
     B("Then: final vector x **LM-head** -> logits -> sample the next token."),
     B("**Tally, one token:** read ~15GB (every weight, once) + KV cache; ~2 x params ≈ **60 GFLOP**."),
     B1("Memory 15GB ÷ 819 GB/s ≈ **18ms**; cores ≈ **2ms** of math -> idle ~8x longer than computing -> **~55 tok/s**, set by the bus."),
     B("Every weight: read once, one multiply-add, discarded. At batch 1 nothing to reuse -> memory-bound. Attention is a small slice; FFN dominates the bytes. (q4 reads half of q8 -> ~2x the q8 ceiling.)")],
    "[1:15]\nSAY: \"Here's that 'read every weight once' claim made concrete. A thirty-billion model at four-bit is about fifteen gigs of weights in the Mac's unified memory, roughly forty-eight layers. To produce a single token the chip walks all forty-eight, once. Each layer, two things. Attention multiplies the current token's vector by the query, key, value, and output matrices, and reads the KV cache, the saved keys and values for every previous token, so it never recomputes the past. Then the feed-forward block, two big matrices, two-thirds of the weights. Across the token the cores stream all fifteen gigs through once, do about sixty billion multiply-adds, and the kicker: the reading takes about eighteen milliseconds, the math about two. The cores sit idle several times longer than they work. Memory-bound, about fifty-five tokens a second, set by the bus.\"\n\nCUE: The 'show your work' slide. Walk the loop once, hammer the 18ms-vs-2ms split. Weights have no reuse at batch 1, so every byte read is on the critical path.",
    font=14)

bullets_slide("Maral: a spare 16GB M2 Air, punching above its weight",
    [B("qwen3:8b / :14b via Ollama's Anthropic endpoint, wired into Claude Code with one env block"),
     B("**Why it clears the Pi:** the M2's unified memory runs **~100 GB/s, ~6x the Pi's ~17** (same theory as the last slides), and 16GB shared RAM (vs 8GB) actually holds an 8-14B model"),
     B("Plot twists:"),
     B1("Tool-use was already at parity with cloud."),
     B1("Real pain wasn't the model, it was memory bandwidth and the WiFi driver crashing under load (rude)"),
     B("16GB is the floor, not a platform. Ran on vibes and about 5 hours of uptime."),
     B("**Takeaway:** the model was never the weak link, the hardware was. 16GB is enough to prove the idea, not to live on it.")],
    "[1:30]\nSAY: \"A spare sixteen-gig M2 MacBook Air carried the whole proof of concept, and it only ran about five hours total. Why it cleared the bar the Pi couldn't: the M2's memory bus is about six times faster, a hundred gigabytes a second versus seventeen, and sixteen gigs of unified RAM actually fits a useful model. The big surprise: tool-use just worked. Claude Code isn't a chatbot, it drives tools through strict function-calling, and my fear was that a small model would faceplant on the mechanics, malformed JSON, the wrong tool. It didn't, it was at parity with cloud out of the box. The real pain was never the model, it was slow memory and the WiFi driver crashing under pressure.\"\n\nCUE: Be precise, it nailed the MECHANICS of tool-calling, separate from driving a long loop to the finish, that convergence problem is Part Three. Takeaway: sixteen gigs proves the idea, can't host it. Sets up 'what do I buy?'")

bullets_slide("The single best tuning knob: kill the thinking tax",
    [B("Qwen3 emits a <think> trace before every answer."),
     B("A coding answer needing 80 tokens cost 600, 8x the work."),
     B("`CLAUDE_CODE_DISABLE_THINKING=1`   (or {\"think\": false})"),
     B("**~3x speedup from one env var.** Everything else I tuned for a week? Secondary.")],
    "[1:30]\nSAY: \"If you remember one config line, it's this. Qwen3 is a reasoning model, it writes a hidden think-trace before every answer. Great for chat, pure overhead for an agent doing lots of tiny tool calls, an eighty-token edit was costing six hundred. One environment variable turns it off, and you get about a three-times real-world speedup.\"\n\nCUE: Audience beat, 'guess how much all my fancy tuning bought me, versus this one line.' Everything else was a rounding error.\n\nALSO TRIED (if asked 'what else did you tune?' - the rest of the ladder, all $0, all secondary):\n- KV-cache quant: OLLAMA_KV_CACHE_TYPE=q8_0 + flash-attn -> 16k-token cache 1.2 GiB (vs ~2.4 at f16), stays 100% on GPU.\n- Context window: OLLAMA_CONTEXT_LENGTH 4k -> 16k -> 4x usable context, still fits memory.\n- Quant sweep: qwen3:8b q4 vs q8 vs 14b-q4 -> q4 won for routine (17.7 tok/s, tied on correctness; bigger/higher-precision bought no quality on bounded tasks, only cost speed).\n- Prompt slim: skillOverrides + ENABLE_TOOL_SEARCH=auto:5 -> Claude Code system prompt 28k -> 9.7k tokens on a cold turn.\n- Keep it warm: OLLAMA_KEEP_ALIVE=5-10m, MAX_LOADED_MODELS, NUM_PARALLEL=1 (avoid cold reloads mid-session).\n- Three-layer think-kill: think:false (Ollama native), thinking:disabled (Anthropic shim), CLAUDE_CODE_DISABLE_THINKING=1 (Claude Code) - NOT on claude-cloud, Opus thinking is worth paying for.\nEach shaved a little; think:false dwarfed the lot combined.")

section_slide("PART TWO", "Reality check & the hardware call",
    notes="[0:05]\nSAY: \"Part two. I've got a working setup, now the uncomfortable part: how good can local actually get, and can you just buy your way to the top? Short answer, no.\"\n\nCUE: Quick beat.")

_rc = table_slide("Reality check: you can't buy frontier parity",
    ["Tier", "Best open weight", "SWE-bench Verified", "Gap to Opus 4.8"],
    [["96GB Mac (what you'd buy)", "Qwen3.6-27B (what I benched)", "77.2%", "−11"],
     ["**Cloud-scale only***", "**DeepSeek-V4 (best open weight, 1.6T)**", "**80.6%**", "**−8 (still short!)**"],
     ["Cloud", "**Opus 4.8**", "**88.6%**", "0"]],
    "**SWE-bench Verified:** 500 real GitHub issues; the score is the % whose generated patch makes the repo's hidden tests pass (functional correctness, not speed). Even the best open weight won't fit the biggest Mac, and it's 8 points behind Opus. **The only thing that gives you Opus quality is Opus.** (llm-stats.com; same 77.2 / 88.6 in Alex Ellis, 'Local AI is not Opus')",
    "[2:00]\nSAY: \"First, what this benchmark even is. SWE-bench Verified is five hundred real GitHub issues from popular open-source projects, and the score is the percentage where the model's patch actually makes the project's hidden test suite pass. So it's pure functional correctness, did it fix the bug, no points for speed. Now the myth-buster: just buy a big enough Mac and you'll match Opus, false. The model I run, Qwen 3.6 27B, resolves seventy-seven percent, eleven behind Opus's eighty-nine. The best open weight on earth, DeepSeek V4, only reaches eighty-point-six, and it won't even fit the biggest Mac Apple sells. You can't spend your way to parity.\"\n\nCUE: Define SWE-bench before the numbers, or the %s mean nothing. Key point: this is CAPABILITY (can it fix it), separate from the wall-clock numbers later. Footnote: not a shill, trying to fire my own bill.",
    col_w=[1.7, 3.4, 2.1, 2.3], font=13)
_rcf = _rc.shapes.add_textbox(Inches(0.6), Inches(6.12), SW - Inches(1.2), Inches(0.8))
_no_autofit(_rcf.text_frame)
_rcf.text_frame.word_wrap = True
def _fn(p, txt):
    r = p.add_run(); r.text = txt
    r.font.size = Pt(9); r.font.italic = True; r.font.color.rgb = DK2
_fn(_rcf.text_frame.paragraphs[0],
    "* DeepSeek-V4 is a 1.6T-param MoE: it won't fit a 512GB Mac, and Apple pulled the 512GB config in March 2026 anyway. Even granting cloud-scale hardware, the best open weight is still −8.")
_fn(_rcf.text_frame.add_paragraph(),
    "** Not an Anthropic shill, I'm trying to fire my own $200/mo Claude bill. The numbers are just what they are.")

_buy = bullets_slide("The buy: a used M3 Ultra 96GB",
    [B("Apple-direct was 4 months backordered (M3 Ultra was EOL'd)"),
     B("Every reseller went dry within days, only the grey market left"),
     B("Verified a sealed eBay unit (buyer protection, 99.3% seller). Wiring four figures to a stranger on eBay is its own personality test."),
     B("**Takeaway:** a just-discontinued machine vanishes from every channel at once. Buy current-gen before the refresh rumor.")],
    "[1:00]\nSAY: \"Right as I decided to buy, Apple discontinued the M3 Ultra, so Apple-direct was four months backordered and every reseller drained within days. I ended up on the grey market, a sealed unit on eBay, and I'll be honest, wiring four figures to a stranger on eBay is its own little personality test.\"\n\nCUE: Breather, tell it like a story. Gesture at the photo, that's the box running the LLM. Takeaway: a just-discontinued machine vanishes from every channel at once.")
# narrow the bullets and drop the Mac Studio photo on the right (the box in question)
_buy.placeholders[1].width = Inches(6.3)
_mac = REPO / "viz" / ("mac-studio-ebay.jpg" if (REPO / "viz" / "mac-studio-ebay.jpg").exists() else "mac-studio.jpg")
_buy.shapes.add_picture(str(_mac), Inches(7.15), Inches(2.4), width=Inches(4.4))
_cap = _buy.shapes.add_textbox(Inches(7.15), Inches(4.85), Inches(4.4), Inches(0.4))
_cr = _cap.text_frame.paragraphs[0].add_run()
_cr.text = "the box running the LLM"; _cr.font.size = Pt(12); _cr.font.italic = True; _cr.font.color.rgb = DK2

section_slide("PART THREE", "The measured verdict",
    notes="[0:05]\nSAY: \"Part three. Enough story, enough vibes. Here are the actual benchmarks, and this is where it surprised me.\"\n\nCUE: Build energy here, this is the best part.")

table_slide("The verdict: local ties cloud on coding",
    ["Model", "Score", "Speed"],
    [["**qwen3-coder:30b (local)**", "**24 / 24**", "**68 tok/s**"],
     ["Opus 4.8 (cloud)", "24 / 24", "-"],
     ["qwen3:32b dense", "16 / 18", "20 tok/s (skip)"]],
    "Mini-bench: 24 algorithmic problems, easy to LeetCode-hard, deterministic pytest scoring. Local tied Opus on every one. **But both went 24/24 - a bench your best model can't lose has stopped measuring. So I built harder tests.**",
    "[1:30]\nSAY: \"Twenty-four coding problems, easy up to LeetCode-hard, scored deterministically with pytest, no model judging itself. The local thirty-billion coder tied Opus on every single one, at sixty-eight tokens a second, for free. But here's the catch: both went twenty-four for twenty-four, so this bench has stopped measuring anything. A benchmark your best model can't lose on is dead weight. What it can't see is the axis that bites day to day, long multi-file agentic work. So I built harder tests.\"\n\nCUE: Stress the rigor, then pivot to the honesty: the bench saturated, so the rest of Part Three is the harder tests. Don't oversell the tie.",
    col_w=[4, 2, 2.5])

table_slide("Why the Studio is fast: bandwidth, not parameters",
    ["Box", "Memory bandwidth", "coder-30b speed"],
    [["MacBook Air M2 16GB", "~100 GB/s", "16 tok/s"],
     ["**Mac Studio M3 Ultra 96GB**", "**819 GB/s**", "**68 tok/s**"]],
    "Same model, 4x faster, purely the memory bus. **MoE (mixture of experts) beats dense:** the dense 32B was mediocre and slow; the 30B MoE activates only 3B params per token, so it streams far less from memory each step, running 3x faster than the dense 32B and scoring higher. (M3 Ultra 819 GB/s per Apple spec; Pi 5 ~17 GB/s.)",
    "[1:30]\nSAY: \"Same model, same quant, eight times the memory bandwidth gives you about four times the tokens per second. The Mac Studio's memory bus is the whole story, it's not about raw compute. And the mixture-of-experts punchline: a thirty-billion MoE that only activates three billion parameters per token beats a dense thirty-two-billion, faster and higher-scoring, because only the active experts stream from memory each token.\"\n\nCUE: The one technical slide. Buying advice: optimize for memory bandwidth, run MoE, don't chase GPU teraflops.",
    col_w=[3.2, 3, 3])

table_slide("But the agent loop is where local gets cooked",
    ["", "pass", "avg wall-clock", "turns (range)"],
    [["**Cloud (Opus)**", "**5 / 5**", "**31 s**", "5-9 (stable)"],
     ["**Local (coder-30b)**", "4 / 5", "**248 s**", "**4-41 (wild)**"]],
    "5 multi-file bug-fixes through the real Claude Code agent loop. **Local 8x slower, a one-line `>` to `>=` fix took it 41 turns / 584 s.**",
    "[2:00]\nSAY: \"The one-shot benchmarks said tie. But I wanted to measure what I actually feel day to day, so I drove the real Claude Code agent loop on five multi-file bug fixes, local versus Opus. Cloud, five out of five, about thirty-one seconds, rock-steady. Local, four out of five, but two hundred forty-eight seconds, eight times slower, and wildly unstable, four to forty-one turns. One of these was a one-line fix, a greater-than to a greater-than-or-equal, and local took forty-one turns and almost ten minutes flailing on it.\"\n\nCUE: THE turning point, give it room. Pause after the 41-turn line, let it hang. The lesson: measure the agent loop, not tokens per second.",
    col_w=[3.4, 1.5, 2.6, 2.6])

image_slide("Watch it: the agent loop side by side",
    REPO / "viz" / "agent-race.gif", width_in=10.2,
    sub="Recording of the real runs at 8x speed, left flails to 9 turns/110s, right is clean at 5 turns/17s.",
    notes="[1:00]\nSAY: \"Same one-line bug. The cloud agent, on the right, has been done for ninety seconds while the local one, on the left, is still grinding.\"\n\nCUE: Let it play silent the first few seconds. Point at the moment the right side freezes green. This is a recording of the real runs, sped up eight times, not a mockup.")

table_slide("Tested it: the instability was the model",
    ["same agent loop, same box", "pass", "turns", "avg time"],
    [["coder-30b (the deck's baseline)", "4 / 5", "**4-41 (wild)**", "248 s"],
     ["Qwen 3.6 27b (q8)", "5 / 5", "7-9", "161 s"],
     ["**gpt-oss 20b**", "**5 / 5**", "**7-9 (stable)**", "**48 s**"],
     ["gemma4 26b", "5 / 5", "8-20", "69 s"],
     ["Opus 4.8 (cloud)", "5 / 5", "5-9", "31 s"]],
    "Every current model is stable and 5/5, the 41-turn coder spiral was a stale model, not local inference. **gpt-oss 20b lands within ~1.5x of cloud**, stable, on a 20B model. The fix was a newer model, not a bigger box.",
    "[1:30]\nSAY: \"Hypothesis from Boykis and Hacker News: maybe that instability is the model, not local inference. So I tested it, same box, same five tasks, swapped my last-gen coder for Qwen three-six, same quant, only the model changed. Five out of five, every task seven to nine turns, dead stable, like cloud. The forty-one-turn spiral became seven. The fix was a newer model, not a bigger box. And gpt-oss 20b lands within about one-and-a-half times cloud.\"\n\nCUE: The payoff: maybe it's software, not hardware. Honest caveat, still slower than cloud, but stable and correct. Most important update since I built the deck.",
    col_w=[3.6, 1.1, 1.8, 1.6], font=12)

table_slide("But on open-ended building, the gap collapses",
    ["backend", "playable", "build time", "turns", "cost"],
    [["Opus 4.8 (cloud)", "**7 / 7**", "50 s", "2", "$0.40"],
     ["qwen3-next:80b (local)", "**7 / 7**", "98 s", "2", "**$0**"],
     ["qwen3-coder:30b (local)", "6 / 7", "166 s", "2", "**$0**"]],
    "Task: 'build playable Space Invaders as one index.html', scored by a Playwright 7-check rubric. **No failing test to spiral on, so local 80B ties cloud, within 2x on speed.**",
    "[1:30]\nSAY: \"Remember the eight-times tax was specifically about debugging, chasing a failing test, where instability compounds. So I tried the opposite: build a playable Space Invaders in one HTML file, scored by an automated browser rubric. On a build-from-scratch task the gap nearly closes, the local eighty-billion gets a perfect seven out of seven, same as Opus, at twice the wall-clock and zero dollars.\"\n\nCUE: The most hopeful data, lift the energy. The bigger local model beat the smaller one on speed AND quality AND tokens. Takeaway: route by task type, build is great for local, debug-a-repo is where cloud earns its keep.",
    col_w=[3.4, 1.6, 1.8, 1.2, 1.4])

image_slide("Same task, three models: all run, only one looks right",
    REPO / "viz" / "appbench" / "side_by_side_2x.gif", width_in=11.2,
    sub="All three are functionally playable (move, shoot, no crash) - but fidelity is a separate axis. coder-30b drew the wrong invader sprites (and hid its game behind a START menu); next-80b rendered invaders as bare rectangles; only Opus 4.8 matched the real Space Invaders look.",
    notes="[1:00]\nSAY: \"These are the actual games the three agents built, played by a script. Let me be honest about 'playable.' All three run, you move, shoot, nothing crashes, that's the functional rubric, and the local ones cost zero. But fidelity is a separate axis, and it's where cloud still wins. The thirty-billion coder drew the wrong invader icons and hid the game behind a start menu; the eighty-billion rendered the invaders as plain rectangles; only Opus four-eight gave a high-fidelity match to the original. So local ties on the logic, cloud still wins the polish, even on a clean build.\"\n\nCUE: Don't oversell 'all playable' - say out loud that only Opus actually looks like Space Invaders. The honest split: function is basically solved locally, asset/visual fidelity isn't yet. Let it loop, invite clone-and-play.")

table_slide("Concrete: real DSG tasks, local vs cloud wall-clock",
    ["Task (internal-web-app stack)", "Type", "Cloud", "Local", "Slower"],
    [["Write an ECR Terraform module", "bounded", "**19 s**", "**186 s**", "**~10x**"],
     ["Scaffold an ECS Fargate service module", "build", "**36 s**", "**256 s**", "**~7x**"],
     ["Fix an ALB OIDC header parser (failing test)", "debug", "**29 s**", "**154 s**", "**~5x**"]],
    "Local got the right answer on every run, same capability, but you wait ~5-10x longer. The gap shrinks as tasks get bigger: bounded is worst (~10x) because cloud finishes it in ~19s, so local's per-turn latency is a bigger multiple. Mean of 3 runs each (warm); local passed all 9. qwen3.6:27b (q8) on the Studio vs Opus 4.8, sandboxed (throwaway dirs, AWS creds stripped, nothing deployed).",
    "[1:30]\nSAY: \"Three real tasks shaped like my own DSG infra: a bounded ECR Terraform module, a build, scaffolding a whole ECS Fargate service, and a debug, fixing a planted bug with a failing test. Each ran three times, local versus cloud, fully sandboxed. Local got the correct answer on all nine runs, same capability. But you wait, five to ten times longer. And the twist, the bounded task is the worst ratio, about ten-x, because cloud finishes it in twenty seconds and local's per-turn latency is a bigger multiple of a tiny number.\"\n\nCUE: The honest counterweight to tokens-per-second. Means of three runs, low variance. 'Local ties' is about whether the answer is right, not how long you wait. gpt-oss narrows the gap.",
    col_w=[4.6, 1.2, 1.2, 1.2, 1.8], font=12)

table_slide("The reconciliation",
    ["Workload", "Local verdict"],
    [["One-shot bounded coding", "**Ties Opus** on capability; ~10x slower wall-clock, free"],
     ["Multi-turn debug loops", "**Current models stable, 5/5; gpt-oss 20b within ~1.5x of cloud**"],
     ["Open-ended building (from scratch)", "**80B ties cloud, ~2x slower, $0**"],
     ["Open-ended repo work (big context)", "**Cloud still wins, on capability**"]],
    "The 41-turn instability was a stale-model artifact; a current model is stable and 5/5. What's left is wall-clock, and it's **model-dependent**: ~1.5x on gpt-oss 20b, ~5-10x on qwen3.6/coder (my dsgbench numbers used the slower qwen3.6). Plus big-context repo work. **Measure the agent loop, keep your model current, route by task.**",
    "[1:00]\nSAY: \"Everything Part Three measured, in one table. One-shot bounded, ties Opus, fast and free. Multi-turn debug loops, current models are stable and five-for-five, gpt-oss within about one-and-a-half times cloud. Open-ended building from scratch, the local eighty-billion ties cloud at twice the wall-clock and zero dollars. Open-ended repo work across big unfamiliar context, cloud still wins on raw capability. The spine of it: local is genuinely capable on most of what I do, and cloud still earns its keep on the hardest repo work, both true at once.\"\n\nCUE: The part-end synthesis, land it slowly. This is the verdict the whole part built to. Let it sit, then move to the money.",
    col_w=[4, 5])

section_slide("PART FOUR", "Economics & the call",
    notes="[0:05]\nSAY: \"Part four. So what does this actually cost, and what should you do on Monday morning?\"\n\nCUE: Quick beat.")

bullets_slide("The call: economics + the setup",
    [B("Cloud ~$200/mo forever; box = one-time, then **$0/mo**. **Breakeven ~2 years.**"),
     B("Privacy: code never leaves the LAN."),
     B("**Hybrid:** `claude` routes to the local Studio (the 80%, daily coding); `claude-cloud` to Opus (the hard 20%)."),
     B("**Sweet spot:** Mac Studio, **64-96GB**, MoE coder model, a few thousand."),
     B("Supplement, not replacement, hard repo work still routes to cloud.")],
    "[1:30]\nSAY: \"The money, then the call. Cloud is two hundred a month forever; the box is a one-time cost, then zero, breakeven about two years. After that it's pure savings, but only for the slice local handles well, it shrinks the bill, it doesn't zero it. So the whole talk reduces to one architecture: hybrid. Two shell aliases, one routes to the local Studio for the eighty percent, one to Opus for the hard twenty, you pick per task. The buying advice: a single dev wants a Mac Studio, sixty-four to ninety-six gigs, an MoE coder, a few thousand dollars.\"\n\nCUE: The Monday-morning slide. If someone photographs one, it's this. Keep saying 'supplement, not replacement.'")

bullets_slide("Caveats & gotchas",
    [B("'Local SOTA at home' is real, narrowly, bounded coding only. Open-ended engineering, no."),
     B("16GB too small, 96GB+ overkill for one dev. Benchmarks saturate, measure your task mix."),
     B("**Swapping models mid-session = instant `API Error: 400`**, /clear or a fresh session."),
     B("**Low quant (q4) silently breaks tool-calling**, q6 is the agentic sweet spot."),
     B("**Forgetting `think:false`** = a silent ~3-8x token tax."),
     B("**The 'instability' was a stale model:** my year-old qwen3-coder:30b spiralled 4-41 turns on a one-line fix; a current Qwen3.6 / gpt-oss does it in 7. Re-bench monthly."),
     B("Ollama traps: restart kills an in-flight `pull`; `ollama ps` shows empty (check `ps aux | grep llama-server`).")],
    "[1:30]\nSAY: \"Let me name my own weaknesses first. 'Local SOTA at home' is true only if you append 'narrowly, bounded coding'; on open-ended engineering it's not. Sixteen gigs is too small to live on, ninety-six-plus is overkill for one person, and benchmarks saturate, the only one that matters is your task mix. Then the operational potholes: swap the model mid-conversation and you get an instant four-hundred error, just clear the session. Low quant quietly breaks tool-calling, run q6 not q4. Forgetting think-false is a silent tax. And to be specific about the instability from earlier: that was a stale model, my year-old coder-30b spiralling four to forty-one turns on a one-line fix, while a current model like Qwen three-six or gpt-oss does the same task in seven, stable. Re-bench monthly.\"\n\nCUE: Q&A insurance plus the practical warnings in one. None are dealbreakers, they're potholes, now you know where they are.",
    font=15)

bullets_slide("What's next",
    [B("**Auto-router on the Pi:** classify each prompt, send bounded work to the Studio + hard repo work to cloud, per-prompt."),
     B("Cut my agents fully over to local; measure on real tasks (tokens, latency, TTFT, success)"),
     B("Verification-loop scaffolding: guardrails took an 8B from 53% to 99% (Forge, Show HN, news.ycombinator.com/item?id=48192383)"),
     B("Instability may be software, not model-size. The fix might be code, not a $10k box.")],
    "[1:00]\nSAY: \"This is live research, not a post-mortem. The one I'm most excited about: an auto-router running on the Pi, the same Pi that was a dead end in Part One, classifying each prompt and sending bounded work to the local Studio, hard repo work to cloud, automatically, prompt by prompt. Beyond that, cut my agents fully over to local, and add verification-loop scaffolding. Published results show guardrails took an eight-billion model from fifty-three to ninety-nine percent. So the instability might be a software problem, the fix could be code, not a ten-thousand-dollar box.\"\n\nCUE: End on momentum. The router is the callback, the dead-end Pi gets redeemed. The last reframe turns the weakness into something solvable.")

_repo = title_slide("Repo + full write-up",
    "github.com/jconnolly/local-llm-pi5",
    notes="[0:30]\nSAY: \"The win is real, and the caveat is real too. Everything's in the repo, reproducible, the benchmarks are real code you can clone and run. Let's open it up.\"\n\nCUE: Say the thesis slowly. Q&A to expect: why not vLLM/MLX, why not always run the 80B, would guardrails fix it, ROI if already paying cloud. Answers in Parts Three and Four. ~30 min here = nailed the pacing.")
# make the repo URL a real hyperlink
for _r in _repo.placeholders[1].text_frame.paragraphs[0].runs:
    _r.hyperlink.address = "https://github.com/jconnolly/local-llm-pi5"

prs.save(str(OUT))
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
print(f"slides: {len(list(prs.slides))}")
